"""
NEST Project Presentation Generator
Creates a professional PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    DARK_BLUE = RGBColor(30, 58, 138)  # #1e3a8a
    LIGHT_BLUE = RGBColor(59, 130, 246)  # #3b82f6
    GREEN = RGBColor(16, 185, 129)  # #10b981
    GRAY = RGBColor(100, 116, 139)  # #64748b
    WHITE = RGBColor(255, 255, 255)
    PURPLE = RGBColor(102, 126, 234)  # #667eea
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 249, 255)
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "NEST"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Near Easy Shop Tracker"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = GRAY
    
    # Description
    desc_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "E-Commerce Grocery Discovery Platform"
    desc_para = desc_frame.paragraphs[0]
    desc_para.alignment = PP_ALIGN.CENTER
    desc_para.font.size = Pt(24)
    desc_para.font.color.rgb = GRAY
    
    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "College Project Presentation"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(18)
    footer_para.font.color.rgb = GRAY
    
    # Slide 2: Abstract
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = WHITE
    
    # Title
    add_slide_title(slide2, "ABSTRACT", DARK_BLUE)
    
    # Content
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    points = [
        ("Problem:", "Traditional grocery shopping faces challenges - difficulty discovering nearby stores, no real-time product availability, inefficient shop management"),
        ("Solution:", "NEST is a comprehensive digital platform connecting local grocery stores with customers through mobile technology"),
        ("Technology:", "Built with Flutter framework for cross-platform support (Android, iOS, Web)"),
        ("Features:", "GPS-based shop discovery, 1000+ product catalog, real-time inventory management, automated billing, and sales analytics"),
        ("Impact:", "Bridges digital divide for local grocery stores while providing modern shopping conveniences to customers")
    ]
    
    for i, (heading, text) in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{heading} {text}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(12)
        # Make heading bold
        p.runs[0].font.bold = True
    
    # Slide 3: Objectives
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide3, "PROJECT OBJECTIVES", DARK_BLUE)
    
    # Left column - Customer Goals
    left_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    left_tf = left_box.text_frame
    
    p = left_tf.paragraphs[0]
    p.text = "🛍️ Customer Goals"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    customer_goals = [
        "Discover nearby stores using GPS",
        "Browse 1000+ product catalog",
        "Real-time availability tracking",
        "Smart shopping cart system"
    ]
    
    for goal in customer_goals:
        p = left_tf.add_paragraph()
        p.text = f"✓ {goal}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(10)
    
    # Right column - Shop Owner Goals
    right_box = slide3.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5.5))
    right_tf = right_box.text_frame
    
    p = right_tf.paragraphs[0]
    p.text = "🏪 Shop Owner Goals"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    owner_goals = [
        "Digital inventory management",
        "Automated billing system",
        "Sales analytics & insights",
        "Business growth tools"
    ]
    
    for goal in owner_goals:
        p = right_tf.add_paragraph()
        p.text = f"✓ {goal}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(10)
    
    # Slide 4: Backend Technologies
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.background.fill.solid()
    slide4.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide4, "BACKEND TECHNOLOGIES", DARK_BLUE)
    
    backend_techs = [
        ("🔥 Firebase Core", "Version 2.24.2 - Cloud platform initialization and configuration"),
        ("☁️ Cloud Firestore", "NoSQL database with real-time synchronization capabilities"),
        ("📦 Provider Pattern", "State management with reactive updates and dependency injection"),
        ("💾 Local Storage", "Offline-first architecture with in-memory data caching")
    ]
    
    y_pos = 1.8
    for i, (title, desc) in enumerate(backend_techs):
        # Create colored box
        box = slide4.shapes.add_shape(
            1,  # Rectangle
            Inches(1), Inches(y_pos), Inches(8), Inches(0.9)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 249, 255)
        box.line.color.rgb = LIGHT_BLUE
        
        # Add text
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        
        y_pos += 1.1
    
    # Slide 5: Frontend Technologies
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    slide5.background.fill.solid()
    slide5.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide5, "FRONTEND TECHNOLOGIES", DARK_BLUE)
    
    frontend_techs = [
        ("📱 Flutter 3.0+", "Cross-platform UI framework for beautiful native apps"),
        ("🎯 Dart 3.0+", "Modern programming language optimized for UI development"),
        ("🎨 Material Design", "Comprehensive UI component library following Google's design system"),
        ("⚡ Provider 6.0", "Reactive state management for seamless UI updates")
    ]
    
    # Grid layout - 2x2
    positions = [
        (Inches(0.5), Inches(1.8)),
        (Inches(5.2), Inches(1.8)),
        (Inches(0.5), Inches(3.8)),
        (Inches(5.2), Inches(3.8))
    ]
    
    for i, ((title, desc), (x, y)) in enumerate(zip(frontend_techs, positions)):
        box = slide5.shapes.add_shape(1, x, y, Inches(4.3), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 249, 255)
        box.line.color.rgb = LIGHT_BLUE
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
    
    # Stats
    stats = [("15+", "Screens"), ("50+", "Components"), ("5000+", "Lines of Code"), ("3", "Platforms")]
    x_pos = 0.8
    for number, label in stats:
        box = slide5.shapes.add_shape(1, Inches(x_pos), Inches(5.8), Inches(2), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = GREEN
        box.line.width = 0
        
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = number
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        
        x_pos += 2.2
    
    # Slide 6: Methodology
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide6, "DEVELOPMENT METHODOLOGY", DARK_BLUE)
    
    # Subtitle
    subtitle_box = slide6.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Agile Scrum - 8 Week Development Sprint"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    
    phases = [
        ("📋 Week 1: Requirements Analysis", "Problem identification, market research, feature prioritization"),
        ("🎨 Week 2: Design", "UI/UX wireframing, database schema, architecture design"),
        ("💻 Week 3-4: Core Development", "Authentication, customer module, map integration, cart system"),
        ("🏪 Week 5-6: Shop Module", "Dashboard, inventory management, billing, analytics"),
        ("✅ Week 7: Testing & QA", "Unit testing, integration testing, bug fixes, optimization"),
        ("🚀 Week 8: Deployment", "APK building, documentation, user manual, release")
    ]
    
    y_pos = 2.1
    for title, desc in phases:
        box = slide6.shapes.add_shape(1, Inches(1.5), Inches(y_pos), Inches(7), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(220, 252, 231)
        box.line.color.rgb = GREEN
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(6, 95, 70)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(4, 120, 87)
        
        y_pos += 0.82
    
    # Slide 7: Algorithms
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    slide7.background.fill.solid()
    slide7.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide7, "KEY ALGORITHMS IMPLEMENTED", DARK_BLUE)
    
    algorithms = [
        ("📍 Haversine Formula", "Distance Calculation", "O(1)"),
        ("🔍 Fuzzy Search", "Multi-field Matching", "O(n×m)"),
        ("📝 Levenshtein Distance", "Edit Distance", "O(n×m)"),
        ("📦 Inventory Management", "Stock Tracking", "O(n)"),
        ("📊 Sales Analytics", "Time-series", "O(n log n)"),
        ("⚡ QuickSort", "Distance Sorting", "O(n log n)")
    ]
    
    # Grid 2x3
    positions = [
        (Inches(0.5), Inches(2)),
        (Inches(3.6), Inches(2)),
        (Inches(6.7), Inches(2)),
        (Inches(0.5), Inches(4.5)),
        (Inches(3.6), Inches(4.5)),
        (Inches(6.7), Inches(4.5))
    ]
    
    for (name, desc, complexity), (x, y) in zip(algorithms, positions):
        box = slide7.shapes.add_shape(1, x, y, Inches(2.9), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(248, 250, 252)
        box.line.color.rgb = RGBColor(226, 232, 240)
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(6)
        
        p3 = tf.add_paragraph()
        p3.text = complexity
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = LIGHT_BLUE
        p3.space_before = Pt(8)
    
    # Slide 8: Expected Output
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide8, "EXPECTED OUTPUT & RESULTS", DARK_BLUE)
    
    # Left - Customer
    left_box = slide8.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(3))
    left_tf = left_box.text_frame
    
    p = left_tf.paragraphs[0]
    p.text = "🛍️ Customer Experience"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    customer_features = [
        "Interactive map with nearby shops",
        "Product search with auto-suggestions",
        "Shopping cart with tax calculations",
        "Real-time stock availability"
    ]
    
    for feature in customer_features:
        p = left_tf.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(8)
    
    # Right - Shop Owner
    right_box = slide8.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(3))
    right_tf = right_box.text_frame
    
    p = right_tf.paragraphs[0]
    p.text = "🏪 Shop Owner Dashboard"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    owner_features = [
        "Daily sales overview with stats",
        "Inventory management (245+ items)",
        "Automated billing system",
        "Analytics with visual charts"
    ]
    
    for feature in owner_features:
        p = right_tf.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(8)
    
    # Stats at bottom
    stats = [("1000+", "Products"), ("GPS", "Tracking"), ("Auto", "Billing"), ("Charts", "Analytics")]
    colors = [LIGHT_BLUE, PURPLE, RGBColor(236, 72, 153), RGBColor(245, 158, 11)]
    x_pos = 1.3
    for (number, label), color in zip(stats, colors):
        box = slide8.shapes.add_shape(1, Inches(x_pos), Inches(5.3), Inches(1.8), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.width = 0
        
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = number
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE
        
        x_pos += 2
    
    # Slide 9: Future Enhancements
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    slide9.background.fill.solid()
    slide9.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide9, "FUTURE ENHANCEMENTS", DARK_BLUE)
    
    enhancements = [
        ("☁️ Phase 2: Backend Integration", "• Firebase real-time sync\n• OTP authentication\n• Multi-device support"),
        ("⭐ Phase 3: Advanced Features", "• Payment gateway (Razorpay/Stripe)\n• Order tracking system\n• Push notifications"),
        ("🤖 Phase 4: AI & ML", "• Product recommendations\n• Demand forecasting\n• Dynamic price optimization"),
        ("📈 Phase 5: Business Expansion", "• Multi-language support\n• Subscription plans\n• Loyalty & rewards program")
    ]
    
    # 2x2 grid
    positions = [
        (Inches(0.7), Inches(2)),
        (Inches(5.3), Inches(2)),
        (Inches(0.7), Inches(4.5)),
        (Inches(5.3), Inches(4.5))
    ]
    
    for (title, desc), (x, y) in zip(enhancements, positions):
        box = slide9.shapes.add_shape(1, x, y, Inches(4.2), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(219, 234, 254)
        box.line.color.rgb = LIGHT_BLUE
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(8)
    
    # Slide 10: Conclusion
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    slide10.background.fill.solid()
    slide10.background.fill.fore_color.rgb = WHITE
    
    add_slide_title(slide10, "CONCLUSION", DARK_BLUE)
    
    # Success heading
    success_box = slide10.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.6))
    tf = success_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ Successfully Delivered"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.color.rgb = GREEN
    p.font.bold = True
    
    # Achievements
    achievements_box = slide10.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(7), Inches(3))
    tf = achievements_box.text_frame
    
    achievements = [
        "Dual-role platform (Customer & Shop Owner)",
        "1000+ product catalog integration",
        "GPS-based shop discovery system",
        "Complete inventory & billing management",
        "Comprehensive sales analytics with charts"
    ]
    
    for i, achievement in enumerate(achievements):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓ {achievement}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(12)
    
    # Final stats
    stats = [("15+", "Screens"), ("5000+", "Lines of Code"), ("3", "Platforms"), ("25MB", "APK Size")]
    x_pos = 1.3
    for number, label in stats:
        box = slide10.shapes.add_shape(1, Inches(x_pos), Inches(5.5), Inches(1.8), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = PURPLE
        box.line.width = 0
        
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = number
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = label
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(11)
        p2.font.color.rgb = WHITE
        
        x_pos += 2
    
    # Thank you
    thanks_box = slide10.shapes.add_textbox(Inches(1), Inches(6.7), Inches(8), Inches(0.6))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You! 🙏"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    
    # Save presentation
    prs.save('NEST_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📁 File: NEST_Presentation.pptx")
    print("📊 Total slides: 10")


def add_slide_title(slide, title_text, color):
    """Helper function to add consistent titles to slides"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = color
    
    # Add line under title
    line = slide.shapes.add_connector(1, Inches(2), Inches(1), Inches(8), Inches(1))
    line.line.color.rgb = color
    line.line.width = Pt(3)


if __name__ == "__main__":
    print("🎨 Creating NEST PowerPoint Presentation...")
    print("=" * 50)
    create_presentation()
    print("=" * 50)
    print("\n📖 Open NEST_Presentation.pptx to view your presentation!")
